"""Extract all text content from PowerPoint files for summarization.

Usage:
    python scripts/extract_pptx_text.py <pptx_file_or_folder> <output_folder>
"""

import sys
from pathlib import Path

from pptx import Presentation


def extract_text_from_pptx(pptx_path: Path) -> str:
    """Extract all text from a .pptx file, slide by slide."""
    prs = Presentation(pptx_path)
    lines: list[str] = []
    lines.append(f"=== FILE: {pptx_path.name} ===")
    lines.append(f"Total slides: {len(prs.slides)}")
    lines.append("")

    for idx, slide in enumerate(prs.slides, start=1):
        lines.append(f"--- Slide {idx} ---")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in para.runs).strip()
                    if text:
                        indent = "  " * (para.level or 0)
                        lines.append(f"{indent}{text}")
            if shape.has_table:
                tbl = shape.table
                lines.append("[TABLE]")
                for row in tbl.rows:
                    cells = [c.text.strip() for c in row.cells]
                    lines.append(" | ".join(cells))
                lines.append("[/TABLE]")
        # Speaker notes
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                lines.append(f"[NOTES] {notes}")
        lines.append("")

    return "\n".join(lines)


def main() -> None:
    if len(sys.argv) != 3:
        print(__doc__)
        sys.exit(1)

    src = Path(sys.argv[1])
    out_dir = Path(sys.argv[2])
    out_dir.mkdir(parents=True, exist_ok=True)

    pptx_files = (
        [src] if src.is_file() else sorted(src.glob("*.pptx"))
    )
    for pptx_file in pptx_files:
        if not pptx_file.exists():
            print(f"Skipping missing file: {pptx_file}")
            continue
        try:
            content = extract_text_from_pptx(pptx_file)
        except Exception as exc:  # noqa: BLE001
            print(f"ERROR processing {pptx_file.name}: {exc}")
            continue
        out_file = out_dir / f"{pptx_file.stem}.txt"
        out_file.write_text(content, encoding="utf-8")
        print(f"Extracted: {pptx_file.name} -> {out_file}")


if __name__ == "__main__":
    main()
